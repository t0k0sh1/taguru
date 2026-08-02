"""The PPTX connector (ADR 0007 §5/§7/§8/§10, issue #352): ``.pptx`` files
into :class:`~taguru_langchain.ingest_connectors.document.ConnectorDocument`,
preserving slide number, slide title, and speaker notes.

``python-pptx`` is kept behind a new optional ``pptx`` extra (ADR 0007
§3/§4's packaging decision, the same one ``pdf.py``/``docx.py`` already
apply to ``pypdf``/``python-docx``) — no dependency forced on a caller who
never ingests PPTX.

A slide's shapes are walked in ``slide.shapes``'s own order (the slide
part's XML document order — the same "trust the format's own structure,
never re-sort by position" posture ``docx.py``'s ``iter_inner_content()``
walk already takes), recursing into a group shape's own nested shapes. Each
non-empty paragraph in a text-frame shape becomes its own paragraph in
``text`` — the same one-``<a:p>``-per-paragraph granularity ``docx.py``
keeps for ``<w:p>`` — and a table becomes exactly ONE paragraph, its rows
joined with ``\\n``, cells with ``" | "``, the same convention ``docx.py``
already uses for a ``<w:tbl>``.

ADR 0007 §7.2 allows at most one locator per paragraph; this connector
spends that budget distinguishing a slide's own body (including its
table(s)) from its speaker notes, rather than — as ``docx.py`` does for its
own, page-less format — distinguishing a table from ordinary body text.
Every body paragraph on slide N (an ordinary text-frame paragraph or a
table) carries ``{"kind": "slide", "value": "N"}``; a slide's speaker notes
are read as their own paragraph(s), each carrying
``{"kind": "speaker_notes", "value": "N"}`` (ADR 0007 §7.3: "never folded
into the paragraph next to it") — appended immediately after that slide's
own body paragraphs, never interleaved with the next slide's. A slide's
title placeholder is read like any other text-frame shape (its text stays
an ordinary body paragraph, carrying the same ``slide`` locator as
everything else on that slide) and, additionally, becomes a
:class:`~.document.SectionEntry` anchored at its own first paragraph — the
same "prose heading, unchanged meaning" ADR 0007 §7.1 keeps `section` to.

No OCR engine, no rasterizer: a presentation left with no extractable text
after the walk (e.g. an image-only deck) degrades to `ocr_required` with
empty text — the same code, and the same "not literally a scan" posture,
``docx.py``/``html.py`` already apply to their own all-image/
boilerplate-only documents. Wiring an external OCR adapter (§10, issue
#352's own other half) to THIS connector is left for a future issue —
:mod:`.ocr` defines the interface generically enough (keyed on
:class:`~taguru.Locator`, not on any one format) for a future connector to
adopt it without a shape change; only ``pdf.py`` calls out to one today.
Encrypted (password-protected) PPTX files use the same MS-OFFCRYPTO
OLE2/CFB container every other Office format does, recognized by its own
magic bytes and reported `encrypted` before ever handing the bytes to
``python-pptx``, the same posture ``docx.py`` already takes.

A chart's own data lives in a separate ``chartN.xml`` part, a SmartArt
diagram's own text in a separate diagram data part, and an embedded/linked
OLE object's content is opaque entirely — none of them reachable through
this connector's own shape walk. A slide referencing any of them is named
in a single `partial_extraction` diagnostic rather than silently
short-changed, the same "structurally present but not walked" posture
``docx.py`` already takes for footnote/endnote/comment/text-box content. A
picture's own visual content (including any text baked into the image
itself) is simply not text this connector — or any connector short of an
OCR adapter — can read; a picture-only slide contributes no paragraphs of
its own, the same silent-omission posture ``docx.py`` already takes for an
image.

``.ppt`` (legacy binary) and ``.pptm`` (macro-enabled) are both
`unsupported_format` — only the plain ``.pptx`` OOXML format is read.
"""

from __future__ import annotations

import hashlib
import io
import re
import zipfile
from collections.abc import Callable
from pathlib import Path
from typing import Final

from taguru import Locator

try:
    from pptx import Presentation as _Presentation
    from pptx.presentation import Presentation as _PresentationType
    from pptx.shapes.group import GroupShape as _GroupShape
    from pptx.slide import Slide as _Slide
except ImportError as _python_pptx_import_error:  # pragma: no cover - exercised only
    # without python-pptx installed
    _Presentation = None  # type: ignore[assignment]
    _PresentationType = None  # type: ignore[assignment, misc]
    _GroupShape = None  # type: ignore[assignment, misc]
    _Slide = None  # type: ignore[assignment, misc]
    _PYTHON_PPTX_IMPORT_ERROR: ImportError | None = _python_pptx_import_error
else:
    _PYTHON_PPTX_IMPORT_ERROR = None

from .._extract import MAX_PASSAGE_BYTES
from ._structure import byte_len, sanitize_paragraph_text
from .document import (
    MAX_SECTION_BYTES,
    ConnectorDocument,
    ConnectorMetadata,
    Diagnostic,
    DiagnosticCode,
    FingerprintInputs,
    LocatorEntry,
    SectionEntry,
    options_digest,
)
from .sources import check_source_id, file_source_id

PARSER_NAME: Final = "taguru-pptx-connector"
PARSER_VERSION: Final = "1.0.0"

_SUPPORTED_SUFFIXES: Final = frozenset({".pptx"})
_CONTENT_TYPE: Final = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

# The raw-file cap, checked before any parsing — same two-cap shape
# PdfConnector/DocxConnector use for the same reason: markup/XML overhead
# means a large-but-sparse deck can stay well under the extracted-text cap,
# and a heavily-styled small deck's raw bytes can still exceed a naive
# per-file cap.
_DEFAULT_MAX_FILE_BYTES: Final = 64 * 1024 * 1024

# MS-OFFCRYPTO password-protected Office documents (including .pptx) are
# re-packaged as a whole OLE2/Compound File Binary container, not a zip —
# this is that container format's own magic number, the same one
# `docx.py` already checks for the same reason. Checked BEFORE any zip/XML
# parsing is attempted, so a password-protected file is reported
# `encrypted` rather than falling through to python-pptx's own zip-open
# failure and being misreported as `corrupt`.
_OLE2_MAGIC: Final = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"

# Content this connector's own shape walk cannot reach at all — named in a
# single `partial_extraction` diagnostic rather than silently omitted.
# Every marker is a literal opening-tag search over each `ppt/slides/
# slideN.xml` part's raw bytes (never a full second parse, the same
# convention `docx.py`'s own `_PARTIAL_CONTENT_MARKERS` already uses): a
# chart's own data lives in a separate `chartN.xml` part; a SmartArt
# diagram's own text lives in a separate diagram data part (`<dgm:relIds>`
# is the reference PowerPoint itself leaves in the slide, pointing at it);
# an embedded/linked OLE object's content is opaque to this connector
# entirely.
_PARTIAL_CONTENT_MARKERS: Final = (
    (b"<c:chart", "chart"),
    (b"<dgm:relIds", "SmartArt diagram"),
    (b"<p:oleObj", "embedded object"),
)

_SLIDE_PART_RE: Final = re.compile(r"^ppt/slides/slide\d+\.xml$")

_EMPTY_SHA256: Final = hashlib.sha256(b"").hexdigest()


class _BodyResult:
    """The one pass over every slide builds all three at once — mirrors
    `docx.py`'s own `_BodyResult`, one paragraph-indexed pass producing
    `text` plus paragraph-indexed `locators`/`sections` together, so they
    can never drift out of sync with each other."""

    __slots__ = ("paragraphs", "locators", "sections", "first_title")

    def __init__(self) -> None:
        self.paragraphs: list[str] = []
        self.locators: list[LocatorEntry] = []
        self.sections: list[SectionEntry] = []
        self.first_title: str | None = None


def _handle_table(shape: object, commit: Callable[[str], int | None]) -> None:
    # `shape.table`: python-pptx dispatches a shape's runtime class by its
    # actual content (text, table, picture, chart, ...) — a content
    # placeholder holding a table becomes one of several table-capable
    # classes (`GraphicFrame`, `TablePlaceholder`, `PlaceholderGraphicFrame`
    # depending on how it was authored), none a stable common ancestor a
    # static isinstance check could narrow to. `has_table` is the library's
    # own documented duck-typed guard for exactly this, checked by the
    # caller before this function is ever called.
    table = shape.table  # type: ignore[attr-defined]
    rows_text = []
    for row in table.rows:
        cells = [cell.text for cell in row.cells]
        if any(cell.strip() for cell in cells):
            rows_text.append(" | ".join(cells))
    if rows_text:
        commit("\n".join(rows_text))


def _build_slide(
    slide: _Slide,
    slide_number: int,
    result: _BodyResult,
    *,
    extract_titles: bool,
    extract_speaker_notes: bool,
    extract_tables: bool,
) -> None:
    body_locator = Locator(kind="slide", value=str(slide_number))

    def commit(text: str, locator: Locator) -> int | None:
        clean = sanitize_paragraph_text(text)
        if not clean:
            return None
        index = len(result.paragraphs)
        result.paragraphs.append(clean)
        result.locators.append(LocatorEntry(paragraph=index, locator=locator))
        return index

    try:
        title_shape = slide.shapes.title if extract_titles else None
    except Exception:  # noqa: BLE001 - a malformed layout reference degrades
        # to "no title", never raised — the same posture
        # DocxConnector._document_title takes for a malformed core
        # properties part.
        title_shape = None
    # `slide.shapes.title` and the same placeholder as yielded by iterating
    # `slide.shapes` below are two distinct Python wrapper objects
    # python-pptx constructs independently around the same underlying XML
    # element — `shape is title_shape` would therefore never match. `.element`
    # (a public accessor on every shape) is what's actually stable across
    # both, so title detection below compares that instead.
    title_element = title_shape.element if title_shape is not None else None

    title_index: int | None = None
    title_label: str | None = None

    def handle_shape(shape: object) -> None:
        nonlocal title_index, title_label
        if isinstance(shape, _GroupShape):
            for child in shape.shapes:
                handle_shape(child)
            return
        if extract_tables and getattr(shape, "has_table", False):
            _handle_table(shape, lambda text: commit(text, body_locator))
            return
        if not getattr(shape, "has_text_frame", False):
            return
        is_title = title_element is not None and getattr(shape, "element", None) is title_element
        # `shape.text_frame`: same duck-typed dispatch `_handle_table`'s own
        # comment explains, for the text-carrying case instead of the
        # table-carrying one — `has_text_frame` is the caller's (this
        # function's own) guard for it, checked immediately above.
        text_frame = shape.text_frame  # type: ignore[attr-defined]
        for paragraph in text_frame.paragraphs:
            index = commit(paragraph.text, body_locator)
            if index is None:
                continue
            if is_title and title_index is None:
                label = paragraph.text.strip()
                if label and byte_len(label) <= MAX_SECTION_BYTES:
                    title_index = index
                    title_label = label

    for shape in slide.shapes:
        handle_shape(shape)

    if title_index is not None and title_label is not None:
        result.sections.append(SectionEntry(paragraph=title_index, section=title_label))
        if result.first_title is None:
            result.first_title = title_label

    if not extract_speaker_notes:
        return
    try:
        has_notes = slide.has_notes_slide
    except Exception:  # noqa: BLE001 - a malformed notes relationship
        # degrades to "no notes", never raised.
        has_notes = False
    if not has_notes:
        return
    try:
        notes_text_frame = slide.notes_slide.notes_text_frame
    except Exception:  # noqa: BLE001 - a malformed notes part degrades to
        # "no notes", never raised.
        notes_text_frame = None
    if notes_text_frame is None:
        return
    notes_locator = Locator(kind="speaker_notes", value=str(slide_number))
    for paragraph in notes_text_frame.paragraphs:
        commit(paragraph.text, notes_locator)


def _build_presentation(
    presentation: _PresentationType,
    *,
    extract_titles: bool,
    extract_speaker_notes: bool,
    extract_tables: bool,
) -> _BodyResult:
    result = _BodyResult()
    for index, slide in enumerate(presentation.slides):
        _build_slide(
            slide,
            index + 1,
            result,
            extract_titles=extract_titles,
            extract_speaker_notes=extract_speaker_notes,
            extract_tables=extract_tables,
        )
    return result


def _partial_extraction_message(raw: bytes) -> str | None:
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as archive:
            slide_names = [name for name in archive.namelist() if _SLIDE_PART_RE.match(name)]
            combined = b"".join(archive.read(name) for name in slide_names)
    except (KeyError, zipfile.BadZipFile):
        return None
    kinds = [name for marker, name in _PARTIAL_CONTENT_MARKERS if marker in combined]
    if not kinds:
        return None
    return f"{', '.join(kinds)} content is present but not extracted by this connector"


def _document_title(presentation: _PresentationType, first_slide_title: str | None) -> str | None:
    try:
        title = presentation.core_properties.title
    except Exception:  # noqa: BLE001 - malformed core properties degrades
        # to "no title", never raised — the same posture
        # DocxConnector._document_title takes for a malformed `/Info`.
        title = None
    if isinstance(title, str) and title.strip():
        return title.strip()
    return first_slide_title


class PptxConnector:
    """Reference connector for ``.pptx`` files (issue #352).

    ``extract_titles`` (default ``True``) controls whether a slide's title
    placeholder becomes a `SectionEntry` — the title's own text is always
    kept as an ordinary body paragraph in `text` either way.
    ``extract_speaker_notes`` (default ``True``) controls whether a slide's
    speaker notes are read at all; when ``False``, notes content is
    skipped entirely rather than folded into the slide's own body.
    ``extract_tables`` (default ``True``) controls whether a table's row/
    cell text is extracted as its own paragraph; when ``False``, a table
    contributes nothing (never folded into a neighboring shape's text).
    ``max_file_bytes`` (default 64 MiB) caps the raw file size checked
    before any parsing.
    """

    def __init__(
        self,
        *,
        extract_titles: bool = True,
        extract_speaker_notes: bool = True,
        extract_tables: bool = True,
        max_file_bytes: int = _DEFAULT_MAX_FILE_BYTES,
    ) -> None:
        if _Presentation is None:
            raise ImportError(
                "PptxConnector requires python-pptx — install with pip install "
                '"langchain-taguru[pptx]"'
            ) from _PYTHON_PPTX_IMPORT_ERROR
        self._extract_titles = extract_titles
        self._extract_speaker_notes = extract_speaker_notes
        self._extract_tables = extract_tables
        self._max_file_bytes = max_file_bytes

    @property
    def parser(self) -> str:
        return PARSER_NAME

    @property
    def parser_version(self) -> str:
        return PARSER_VERSION

    @property
    def parse_options_digest(self) -> str:
        """The digest :meth:`read` would stamp into
        ``fingerprint_inputs.parse_options_digest`` for THIS instance's
        effective config, computable without reading anything (ADR 0007
        §6.3, added for issue #351's connector-level checkpoint: it gates
        skipping a re-parse before ever calling :meth:`read`, which needs
        the full candidate fingerprint up front)."""
        return self._options_digest()

    def supports(self, reference: str) -> bool:
        return Path(reference).suffix.lower() in _SUPPORTED_SUFFIXES

    def _options_digest(self) -> str:
        return options_digest(
            {
                "extract_titles": self._extract_titles,
                "extract_speaker_notes": self._extract_speaker_notes,
                "extract_tables": self._extract_tables,
                "max_file_bytes": self._max_file_bytes,
            }
        )

    def _fingerprint(self, raw_content_sha256: str) -> FingerprintInputs:
        return FingerprintInputs(
            raw_content_sha256=raw_content_sha256,
            parser=self.parser,
            parser_version=self.parser_version,
            parse_options_digest=self._options_digest(),
        )

    def _failure(
        self,
        *,
        source: str,
        display_name: str,
        code: DiagnosticCode,
        message: str,
        raw_content_sha256: str,
        content_type: str | None = _CONTENT_TYPE,
    ) -> ConnectorDocument:
        return ConnectorDocument(
            source=source,
            text="",
            metadata=ConnectorMetadata(
                origin_uri=source, display_name=display_name, content_type=content_type
            ),
            fingerprint_inputs=self._fingerprint(raw_content_sha256),
            diagnostics=(Diagnostic(code=code, message=message, source=source),),
        )

    def read(self, reference: str) -> ConnectorDocument:
        # As in TextFileConnector/PdfConnector/DocxConnector: `reference`
        # itself is the source id (ADR 0007 §6.1); `Path` is only ever
        # used for filesystem access below, never for identity.
        path = Path(reference)
        source = file_source_id(reference)
        display_name = path.name

        def failure(
            code: DiagnosticCode,
            message: str,
            raw_content_sha256: str = _EMPTY_SHA256,
            *,
            content_type: str | None = _CONTENT_TYPE,
        ) -> ConnectorDocument:
            return self._failure(
                source=source,
                display_name=display_name,
                code=code,
                message=message,
                raw_content_sha256=raw_content_sha256,
                content_type=content_type,
            )

        source_diagnostic = check_source_id(source)
        if source_diagnostic is not None:
            return failure(source_diagnostic.code, source_diagnostic.message)

        if not self.supports(reference):
            # Unlike every other failure below, the extension mismatch here
            # is affirmative evidence the file is NOT a PPTX — content_type
            # stays unclaimed (None), the same posture DocxConnector's own
            # `read` takes for its own unsupported-format case.
            return failure(
                "unsupported_format",
                f"unsupported extension {path.suffix!r} (only .pptx)",
                content_type=None,
            )

        try:
            size = path.stat().st_size
        except OSError as error:
            return failure("unreadable", str(error))
        if size > self._max_file_bytes:
            return failure(
                "content_too_large",
                f"{size} bytes exceeds the {self._max_file_bytes}-byte file cap",
            )

        try:
            raw = path.read_bytes()
        except OSError as error:
            return failure("unreadable", str(error))
        # The raw bytes, before parsing — ADR 0007 §5's connector-native
        # checkpoint fingerprint, deliberately distinct from any hash of
        # the extracted text (§6.2).
        raw_content_sha256 = hashlib.sha256(raw).hexdigest()

        if raw.startswith(_OLE2_MAGIC):
            return failure(
                "encrypted",
                "the document requires a password this connector does not have",
                raw_content_sha256,
            )

        try:
            presentation = _Presentation(io.BytesIO(raw))
            body = _build_presentation(
                presentation,
                extract_titles=self._extract_titles,
                extract_speaker_notes=self._extract_speaker_notes,
                extract_tables=self._extract_tables,
            )
        except Exception as error:  # noqa: BLE001 - python-pptx raises a
            # grab-bag of zipfile/lxml/KeyError/ValueError for a malformed
            # package, none sharing a common base — every such failure is
            # still this document's own structure failing to parse.
            return failure("corrupt", str(error), raw_content_sha256)

        if not body.paragraphs:
            return failure(
                "ocr_required",
                "no extractable text (an image-only presentation, or a genuinely empty one)",
                raw_content_sha256,
            )

        text = "\n\n".join(body.paragraphs)
        if byte_len(text) > MAX_PASSAGE_BYTES:
            return failure(
                "content_too_large",
                f"{byte_len(text)} bytes of extracted text exceeds the "
                f"{MAX_PASSAGE_BYTES}-byte passage cap",
                raw_content_sha256,
            )

        diagnostics: list[Diagnostic] = []
        partial_message = _partial_extraction_message(raw)
        if partial_message is not None:
            diagnostics.append(
                Diagnostic(code="partial_extraction", message=partial_message, source=source)
            )

        return ConnectorDocument(
            source=source,
            text=text,
            locators=tuple(body.locators),
            sections=tuple(body.sections),
            metadata=ConnectorMetadata(
                origin_uri=source,
                display_name=display_name,
                title=_document_title(presentation, body.first_title),
                content_type=_CONTENT_TYPE,
            ),
            fingerprint_inputs=self._fingerprint(raw_content_sha256),
            diagnostics=tuple(diagnostics),
        )
