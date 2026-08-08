"""Minimal ``.pptx`` byte-builders for
:mod:`taguru_langchain.ingest_connectors.pptx` tests (issue #352).

Unlike ``tests/_docx.py``'s own "hand-assemble a raw zip of exactly the
parts needed" convention, these are built through ``python-pptx``'s real
``Presentation()``/``.save()`` writer. A ``.docx`` package's body lives in
one part (``word/document.xml``) alongside a couple of small, easy-to-hand-
write siblings; a ``.pptx`` package spans several genuinely interdependent
parts even for one blank slide (``ppt/presentation.xml``, a slide master, at
least one slide layout, a theme, and each slide's own relationships part) —
hand-rolling all of them correctly is far more fragile than the payoff, and
every fixture built this way is, for free, real ``python-pptx`` output
(exactly what ``test_docx_connector.py``'s own "a real library writer's
output round-trips" test has to earn separately for DOCX).

Composable primitives (``blank_presentation``, ``add_title_slide``,
``add_body``, ``add_table``, ``add_group``, ``add_notes``, ``add_chart``,
``save_bytes``) let a test build exactly the shape it needs — multiple
slides, mixed content, nested groups — rather than a single rigid one-slide
helper. A SmartArt diagram and an embedded/linked OLE object are the two
exceptions: ``python-pptx`` has no create-API for either, so
``with_smartart``/``with_ole_object`` patch their own reference marker
directly into an already-built deck's raw ``ppt/slides/slideN.xml`` bytes —
the same "the reference alone is enough, its real target part is never
declared" trick ``tests/_docx.py``'s own ``footnote_reference_para`` uses
for content whose target part ``PptxConnector``'s own marker scan never
reads either.
"""

from __future__ import annotations

import io
import struct
import zipfile
import zlib
from collections.abc import Sequence

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pptx.util import Inches

_TITLE_ONLY_LAYOUT = 5
_BLANK_LAYOUT = 6
_TITLE_AND_CONTENT_LAYOUT = 1


def blank_presentation() -> PresentationType:
    return Presentation()


def set_core_title(presentation: PresentationType, title: str) -> None:
    """The presentation-level ``core_properties.title`` — deliberately a
    separate helper from ``add_title_slide``'s own ``title`` parameter,
    which sets a SLIDE's title placeholder text, not this document-level
    property; ``PptxConnector``'s own ``_document_title`` prefers this one,
    the same precedence ``docx_bytes(body, title=...)`` documents for
    ``DocxConnector``."""
    presentation.core_properties.title = title


def add_title_slide(presentation: PresentationType, title: str | None) -> Slide:
    """Adds one slide. ``title=None`` uses the ``Blank`` layout (no title
    placeholder at all — ``slide.shapes.title`` is ``None``); any other
    value (including ``""``) uses ``Title Only`` (a title placeholder is
    present; ``""`` leaves it with no text, the shape a locator-anchor test
    needs to prove an empty title creates no ``SectionEntry``)."""
    layout_index = _BLANK_LAYOUT if title is None else _TITLE_ONLY_LAYOUT
    slide = presentation.slides.add_slide(presentation.slide_layouts[layout_index])
    if title:
        title_shape = slide.shapes.title
        assert title_shape is not None
        title_shape.text_frame.text = title
    return slide


def add_title_and_content_slide(presentation: PresentationType, *, title: str, body: str) -> Slide:
    """A slide built from the ``Title and Content`` layout, its title AND
    body both populated through their own real placeholder shapes — never
    a manually added textbox, unlike every other helper here — proving
    ``PptxConnector`` reads an ordinary content placeholder exactly like
    the plain textboxes every other fixture in this module uses."""
    slide = presentation.slides.add_slide(presentation.slide_layouts[_TITLE_AND_CONTENT_LAYOUT])
    title_shape = slide.shapes.title
    assert title_shape is not None
    title_shape.text_frame.text = title
    body_placeholder = slide.placeholders[1]
    body_placeholder.text_frame.text = body
    return slide


def add_body(slide: Slide, paragraphs: Sequence[str]) -> None:
    """A plain textbox holding one paragraph per ``paragraphs`` element —
    the general body-shape case every ordinary bullet/text box slide
    exercises."""
    box = slide.shapes.add_textbox(Inches(0), Inches(1), Inches(4), Inches(3))
    text_frame = box.text_frame
    text_frame.text = paragraphs[0]
    for extra in paragraphs[1:]:
        text_frame.add_paragraph().text = extra


def add_table(slide: Slide, rows: Sequence[Sequence[str]]) -> None:
    row_count, col_count = len(rows), len(rows[0])
    graphic_frame = slide.shapes.add_table(
        row_count, col_count, Inches(0), Inches(0), Inches(4), Inches(2)
    )
    table = graphic_frame.table
    for row_index, row in enumerate(rows):
        for col_index, cell in enumerate(row):
            table.cell(row_index, col_index).text = cell


def add_group(slide: Slide, texts: Sequence[str]) -> None:
    """A single group shape holding one textbox per ``texts`` element —
    exercises ``PptxConnector``'s own recursion into a group's nested
    shapes."""
    group = slide.shapes.add_group_shape()
    top = 0
    for text in texts:
        box = group.shapes.add_textbox(Inches(0), top, Inches(1), Inches(1))
        box.text_frame.text = text
        top += Inches(1)


def add_notes(slide: Slide, paragraphs: Sequence[str]) -> None:
    text_frame = slide.notes_slide.notes_text_frame
    text_frame.text = paragraphs[0]
    for extra in paragraphs[1:]:
        text_frame.add_paragraph().text = extra


def add_chart(slide: Slide) -> None:
    """A real chart — its own data lives in a separate ``chartN.xml`` part
    ``PptxConnector`` never reads, so this exercises the ``chart`` case of
    its ``partial_extraction`` marker scan without needing any hand-patched
    XML at all."""
    chart_data = CategoryChartData()
    chart_data.categories = ["a", "b"]
    chart_data.add_series("Series 1", (1, 2))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0), Inches(0), Inches(2), Inches(2), chart_data
    )


def save_bytes(presentation: PresentationType) -> bytes:
    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()


def pptx_bytes(
    *,
    title: str | None = None,
    bodies: Sequence[str] = (),
    table: Sequence[Sequence[str]] | None = None,
    notes: Sequence[str] = (),
    chart: bool = False,
) -> bytes:
    """One-slide convenience wrapper over the primitives above, for the
    common case of a single slide exercising one or two features at once."""
    presentation = blank_presentation()
    slide = add_title_slide(presentation, title)
    if bodies:
        add_body(slide, bodies)
    if table is not None:
        add_table(slide, table)
    if notes:
        add_notes(slide, notes)
    if chart:
        add_chart(slide)
    return save_bytes(presentation)


# A SmartArt diagram, minimal enough to parse: a `<p:graphicFrame>` whose
# `graphicData` names the diagram namespace, containing the `<dgm:relIds>`
# reference `PptxConnector`'s own marker scan looks for. `r:dm`/`r:lo`/
# `r:qs`/`r:cs` name relationship ids this deck never actually declares —
# `python-pptx` accepts an unrecognized `graphicData` uri without resolving
# them (confirmed: `has_table`/`has_text_frame`/`shape_type` all degrade to
# False/None rather than raising), so this needs no real diagram data part
# behind it, only the reference PowerPoint itself leaves in the slide.
_SMARTART_MARKER = (
    b'<p:graphicFrame><p:nvGraphicFramePr><p:cNvPr id="9001" name="Diagram"/>'
    b"<p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr>"
    b'<p:xfrm><a:off x="0" y="0"/><a:ext cx="1" cy="1"/></p:xfrm>'
    b'<a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/diagram">'
    b'<dgm:relIds xmlns:dgm="http://schemas.openxmlformats.org/drawingml/2006/diagram" '
    b'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
    b'r:dm="rIdX" r:lo="rIdX" r:qs="rIdX" r:cs="rIdX"/></a:graphicData></a:graphic>'
    b"</p:graphicFrame>"
)

# Likewise minimal for an embedded OLE object: a `<p:graphicFrame>` whose
# `graphicData` names the presentation-ole namespace, containing the
# `<p:oleObj>` reference — the real shape PowerPoint emits wraps this in an
# `mc:AlternateContent` choice/fallback pair, trimmed here to the one
# element `PptxConnector`'s own marker scan actually greps for.
_OLE_OBJECT_MARKER = (
    b'<p:graphicFrame><p:nvGraphicFramePr><p:cNvPr id="9002" name="Object"/>'
    b"<p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr>"
    b'<p:xfrm><a:off x="0" y="0"/><a:ext cx="1" cy="1"/></p:xfrm>'
    b'<a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/presentationml/2006/ole">'
    b'<p:oleObj name="Object" r:id="rIdX" imgW="1" imgH="1" progId="Package" '
    b'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
    b"<p:embed/></p:oleObj></a:graphicData></a:graphic>"
    b"</p:graphicFrame>"
)


def _inject_marker(raw: bytes, marker_xml: bytes, *, slide_number: int) -> bytes:
    name = f"ppt/slides/slide{slide_number}.xml"
    with zipfile.ZipFile(io.BytesIO(raw)) as archive:
        entries = {part: archive.read(part) for part in archive.namelist()}
    xml = entries[name]
    patched = xml.replace(b"</p:spTree>", marker_xml + b"</p:spTree>")
    assert patched != xml, f"marker injection failed: {name} has no </p:spTree>"
    entries[name] = patched
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as archive:
        for part, data in entries.items():
            archive.writestr(part, data)
    return buf.getvalue()


def with_smartart(raw: bytes, *, slide_number: int = 1) -> bytes:
    return _inject_marker(raw, _SMARTART_MARKER, slide_number=slide_number)


def with_ole_object(raw: bytes, *, slide_number: int = 1) -> bytes:
    return _inject_marker(raw, _OLE_OBJECT_MARKER, slide_number=slide_number)


def zip_bomb_pptx(declared_size: int = 2 * 1024 * 1024) -> bytes:
    """A small, high-ratio zip whose one entry declares ``declared_size``
    zero bytes but compresses to almost nothing — the shape
    `PptxConnector`'s `max_decompressed_bytes` cap exists to refuse before
    python-pptx ever decompresses it. Never assembled into a fully valid
    pptx package: the cap is checked, and this document refused, before
    python-pptx would ever open it."""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("ppt/slides/slide1.xml", b"\0" * declared_size)
    return buf.getvalue()


def forged_size_zip_bomb_pptx(payload_size: int = 32 * 1024 * 1024) -> bytes:
    """A zip bomb that DEFEATS a metadata-only check: its one entry really
    decompresses to ``payload_size`` bytes, but every ``file_size`` field
    (local header and central directory) is forged down to 50, so
    ``ZipInfo.file_size`` reports 50 while the deflate stream still expands
    to megabytes — the exact bypass `decompressed_size_within` closes."""
    payload = b"\0" * payload_size
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("ppt/slides/slide1.xml", payload)
    raw = bytearray(buf.getvalue())
    central = raw.find(b"PK\x01\x02")
    local = raw.find(b"PK\x03\x04")
    forged = 50
    struct.pack_into("<I", raw, central + 24, forged)
    struct.pack_into("<I", raw, local + 22, forged)
    crc = zlib.crc32(payload[:forged])
    struct.pack_into("<I", raw, central + 16, crc)
    struct.pack_into("<I", raw, local + 14, crc)
    return bytes(raw)


def corrupt_pptx(kind: str = "not_zip") -> bytes:
    """Three distinct corruption shapes ``PptxConnector`` must all report
    as `corrupt`, never raise: ``"not_zip"`` (not a zip container at all),
    ``"missing_part"`` (a valid zip missing ``ppt/presentation.xml``
    entirely), and ``"malformed_xml"`` (a slide part exists but its XML
    doesn't parse)."""
    if kind == "not_zip":
        return b"this is not a zip file at all, just plain bytes"
    baseline = pptx_bytes(bodies=["placeholder"])
    with zipfile.ZipFile(io.BytesIO(baseline)) as archive:
        entries = {part: archive.read(part) for part in archive.namelist()}
    if kind == "missing_part":
        del entries["ppt/presentation.xml"]
    elif kind == "malformed_xml":
        entries["ppt/slides/slide1.xml"] = b"<p:sld><p:cSld><p:spTree not closed"
    else:
        raise ValueError(f"unknown corrupt_pptx kind: {kind!r}")
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as archive:
        for part, data in entries.items():
            archive.writestr(part, data)
    return buf.getvalue()


def encrypted_pptx() -> bytes:
    """The first 8 bytes any MS-OFFCRYPTO password-protected Office file
    starts with (the OLE2/Compound File Binary container signature) —
    enough for ``PptxConnector``'s own magic-byte check, which never
    parses past this header. Identical to ``tests/_docx.py``'s own
    ``encrypted_docx``: the signature is format-agnostic."""
    return b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 504
